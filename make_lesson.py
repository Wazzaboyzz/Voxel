#!/usr/bin/env python3
"""
make_lesson.py - one command, topic in, ready-to-teach PPT out.

Pipeline:
  1. Nemotron 3 Ultra (free, via OpenRouter) writes the lesson outline/script
     (via content_provider.generate_outline - shared across Voxel)
  2. Pollinations.ai generates one illustration per slide, from the SAME
     outline above - no separate/duplicate outline call
     (via image_provider.generate_all_images - shared across Voxel)
  3. python-pptx builds the .pptx deck from that outline
  4. Piper (free, local TTS) generates narration audio per slide
  5. Freesound.org (free API) pulls background music/SFX matching each slide's mood
  6. ffmpeg stitches narration + music into each slide's embedded audio
  7. A canonical project.json manifest is written for the run (project_provider)

Usage:
    python make_lesson.py "Present perfect tense for intermediate ESL students"

Required environment variables (set these before running):
    OPENROUTER_API_KEY   - your OpenRouter key (free Nemotron access)
    FREESOUND_API_KEY    - free account key from freesound.org/apiv2/apply/

Required local tools (install once, all free):
    pip install python-pptx requests --break-system-packages
    Piper TTS binary: https://github.com/rhasspy/piper (download a voice model too)
    ffmpeg: must be on PATH
"""

import os
import sys
import subprocess
from pathlib import Path

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from content_provider import generate_outline
from image_provider import generate_all_images as generate_all_slide_images, LESSON_STYLE_SUFFIX
from project_provider import build_lesson_project_record, write_project_json


FREESOUND_API_KEY = os.environ.get("FREESOUND_API_KEY", "")
FREESOUND_SEARCH_URL = "https://freesound.org/apiv2/search/text/"

PIPER_BIN = os.environ.get("PIPER_BIN", "piper")
PIPER_VOICE = os.environ.get("PIPER_VOICE", "en_US-lessac-medium.onnx")

OUTPUT_DIR = Path("output")


def build_deck(topic, slides, out_path, image_files=None):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.text = slide_data.get("title", f"Slide {i+1}")
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.2), Inches(4.5))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        bullets = slide_data.get("body", [])
        for j, bullet in enumerate(bullets):
            p = body_tf.paragraphs[0] if j == 0 else body_tf.add_paragraph()
            p.text = f"- {bullet}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        image_file = image_files[i] if image_files else None
        if image_file and Path(image_file).exists():
            slide.shapes.add_picture(str(image_file), Inches(6.0), Inches(1.6), height=Inches(4.5))

        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data.get("narration", "")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return prs


def generate_narration(slides, audio_dir):
    audio_dir.mkdir(parents=True, exist_ok=True)
    narration_files = []
    for i, slide_data in enumerate(slides):
        text = slide_data.get("narration", "").strip()
        out_wav = audio_dir / f"slide_{i+1:02d}_narration.wav"
        if not text:
            narration_files.append(None)
            continue
        try:
            subprocess.run(
                [PIPER_BIN, "--model", PIPER_VOICE, "--output_file", str(out_wav)],
                input=text, text=True, check=True, capture_output=True,
            )
            narration_files.append(out_wav)
        except FileNotFoundError:
            print(f"  [warn] Piper binary not found at '{PIPER_BIN}'. Skipping narration.")
            narration_files.append(None)
        except subprocess.CalledProcessError as e:
            print(f"  [warn] Piper failed on slide {i+1}: {e.stderr}")
            narration_files.append(None)
    return narration_files


def fetch_background_audio(slides, audio_dir):
    audio_dir.mkdir(parents=True, exist_ok=True)
    if not FREESOUND_API_KEY:
        print("  [warn] FREESOUND_API_KEY not set - skipping background music/SFX.")
        return [None] * len(slides)
    bg_files = []
    for i, slide_data in enumerate(slides):
        mood = slide_data.get("mood", "calm neutral")
        out_path = audio_dir / f"slide_{i+1:02d}_background.mp3"
        try:
            resp = requests.get(
                FREESOUND_SEARCH_URL,
                params={"query": mood, "token": FREESOUND_API_KEY,
                        "filter": "duration:[5 TO 60]", "fields": "id,name,previews", "page_size": 1},
                timeout=30,
            )
            resp.raise_for_status()
            results = resp.json().get("results", [])
            if not results:
                bg_files.append(None)
                continue
            preview_url = results[0]["previews"].get("preview-hq-mp3")
            if not preview_url:
                bg_files.append(None)
                continue
            audio_data = requests.get(preview_url, timeout=30)
            audio_data.raise_for_status()
            out_path.write_bytes(audio_data.content)
            bg_files.append(out_path)
        except requests.RequestException as e:
            print(f"  [warn] Freesound fetch failed for slide {i+1} (mood: {mood}): {e}")
            bg_files.append(None)
    return bg_files


def mix_audio(narration_files, bg_files, mix_dir):
    mix_dir.mkdir(parents=True, exist_ok=True)
    mixed_files = []
    for i, (narration, bg) in enumerate(zip(narration_files, bg_files)):
        out_path = mix_dir / f"slide_{i+1:02d}_mixed.mp3"
        if narration is None and bg is None:
            mixed_files.append(None)
            continue
        if narration and bg:
            cmd = ["ffmpeg", "-y", "-i", str(narration), "-i", str(bg), "-filter_complex",
                   "[1:a]volume=0.25[bg];[0:a][bg]amix=inputs=2:duration=first:dropout_transition=2", str(out_path)]
        elif narration:
            cmd = ["ffmpeg", "-y", "-i", str(narration), str(out_path)]
        else:
            cmd = ["ffmpeg", "-y", "-i", str(bg), str(out_path)]
        try:
            subprocess.run(cmd, check=True, capture_output=True)
            mixed_files.append(out_path)
        except FileNotFoundError:
            print("  [warn] ffmpeg not found on PATH - skipping audio mix for this slide.")
            mixed_files.append(None)
        except subprocess.CalledProcessError as e:
            print(f"  [warn] ffmpeg mix failed on slide {i+1}: {e.stderr.decode(errors='ignore')}")
            mixed_files.append(None)
    return mixed_files


def main():
    if len(sys.argv) < 2:
        print('Usage: python make_lesson.py "Your lesson topic here"')
        sys.exit(1)

    topic = " ".join(sys.argv[1:])
    safe_name = "".join(c if c.isalnum() or c in " -_" else "" for c in topic).strip().replace(" ", "_")[:60]
    run_dir = OUTPUT_DIR / safe_name

    print(f"Generating lesson outline for: {topic}")
    slides = generate_outline(topic)
    print(f"  -> {len(slides)} slides planned")

    print("Generating slide images (Pollinations.ai, free, no key required)...")
    image_files = generate_all_slide_images(
        slides, run_dir / "images", filename_prefix="slide",
        width=1024, height=576, style_suffix=LESSON_STYLE_SUFFIX,
    )

    print("Building .pptx deck...")
    deck_path = run_dir / f"{safe_name}.pptx"
    build_deck(topic, slides, deck_path, image_files=image_files)
    print(f"  -> {deck_path}")

    print("Generating narration audio (Piper)...")
    narration_files = generate_narration(slides, run_dir / "narration")

    print("Fetching background music/SFX (Freesound)...")
    bg_files = fetch_background_audio(slides, run_dir / "background")

    print("Mixing final audio per slide (ffmpeg)...")
    mixed_files = mix_audio(narration_files, bg_files, run_dir / "mixed")

    print("Writing project.json manifest...")
    record = build_lesson_project_record(
        topic=topic,
        slides=slides,
        image_files=image_files,
        deck_path=deck_path,
        narration_files=narration_files,
        bg_files=bg_files,
        mixed_files=mixed_files,
        run_dir=run_dir,
    )
    project_json_path = write_project_json(record, run_dir)
    print(f"  -> {project_json_path}")

    print()
    print("Done. Output folder:")
    print(f"  {run_dir.resolve()}")
    print()
    print("NOTE: audio files are generated but not yet embedded into the .pptx")
    print("itself. Open the deck and manually insert each slide's mixed audio")
    print("file via PowerPoint's Insert > Audio menu, matching by slide number.")


if __name__ == "__main__":
    main()


if __name__ == "__main__":
    main()
